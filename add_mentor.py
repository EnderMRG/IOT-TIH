"""
Script to add Ashish Kumar Mahato as a mentor on slide 2 of the PowerPoint.
Copies the same card style as existing team members.
"""
from lxml import etree
from pptx import Presentation

INPUT_PATH = r'i:\Proj\IOT  TIH\report\FloodEye_Final_Conference_Edition.pptx'
OUTPUT_PATH = r'i:\Proj\IOT  TIH\report\FloodEye_Final_Conference_Edition_WithMentor.pptx'

p = Presentation(INPUT_PATH)
slide = p.slides[1]

# -- Slide dimensions
SLIDE_W = p.slide_width   # 12192000
SLIDE_H = p.slide_height  # 6858000

# -- Card dimensions (from existing member cards)
CARD_W           = 2468880
CARD_H           = 3291840
RR_SHADOW_DX     = 45720
RECT_H           = 640080
OVAL_W           = 822960
OVAL_H           = 822960
OVAL_REL_X       = 822960
OVAL_REL_Y       = 320040
DIVIDER_W        = 1371600
DIVIDER_H        = 22860
DIVIDER_REL_X    = 548640

# -- Colors
COLOR_CARD_SHADOW = 'D7E1EB'
COLOR_CARD_BG     = 'FFFFFF'
COLOR_HEADER      = '008296'
COLOR_OVAL_BORDER = '008296'
COLOR_INITIALS    = '008296'
COLOR_NAME        = '232832'
COLOR_DIVIDER     = '008296'
COLOR_ROLE        = '46505F'

# -- Reference vertical positions from existing cards
CARD_SHADOW_TOP = 2057400
CARD_TOP        = 2011680

# -- Mentor card positioning
MENTOR_SECTION_Y       = 4780000   # "Our Mentor" label y
MENTOR_CARD_TOP        = 5010000   # top of mentor card
MENTOR_CARD_SHADOW_TOP = MENTOR_CARD_TOP + (CARD_SHADOW_TOP - CARD_TOP)

MENTOR_CARD_LEFT  = (SLIDE_W - CARD_W) // 2
MENTOR_SHADOW_LEFT = MENTOR_CARD_LEFT + RR_SHADOW_DX

OVAL_LEFT        = MENTOR_CARD_LEFT + OVAL_REL_X
OVAL_TOP         = MENTOR_CARD_TOP  + OVAL_REL_Y

INITIALS_TB_LEFT = OVAL_LEFT
INITIALS_TB_TOP  = MENTOR_CARD_TOP + (2514600 - 2011680)  # 502920 from card top

NAME_TB_LEFT     = MENTOR_CARD_LEFT + (914400 - 822960)   # 91440 from card left
NAME_TB_TOP      = MENTOR_CARD_TOP  + (3383280 - 2011680) # 1371600 from card top

DIVIDER_LEFT     = MENTOR_CARD_LEFT + DIVIDER_REL_X
DIVIDER_TOP      = MENTOR_CARD_TOP  + (3886199 - 2011680) # ~1874519 from card top

ROLE_TB_LEFT     = NAME_TB_LEFT
ROLE_TB_TOP      = MENTOR_CARD_TOP  + (4114800 - 2011680) # 2103120 from card top


def next_id(slide):
    ids = [s.shape_id for s in slide.shapes]
    return max(ids) + 1


def make_sp(shape_id, name, x, y, cx, cy, prst,
            fill_color=None, no_fill=False,
            line_w=None, line_color=None, no_line=False,
            txbody_xml=None, is_textbox=False):
    """Build a <p:sp> element from scratch."""
    cNvSpPr_extra = ' txBox="1"' if is_textbox else ''

    if no_fill:
        fill_xml = '<a:noFill/>'
    elif fill_color:
        fill_xml = f'<a:solidFill><a:srgbClr val="{fill_color}"/></a:solidFill>'
    else:
        fill_xml = ''

    if no_line:
        line_xml = '<a:ln><a:noFill/></a:ln>'
    elif line_color:
        w_attr = f' w="{line_w}"' if line_w else ''
        line_xml = (
            f'<a:ln{w_attr}>'
            f'<a:solidFill><a:srgbClr val="{line_color}"/></a:solidFill>'
            f'</a:ln>'
        )
    else:
        line_xml = ''

    if txbody_xml is None:
        txbody_xml = (
            '<p:txBody><a:bodyPr/><a:lstStyle/>'
            '<a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p></p:txBody>'
        )

    xml_str = (
        f'<p:sp'
        f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        f'<p:nvSpPr>'
        f'<p:cNvPr id="{shape_id}" name="{name}"/>'
        f'<p:cNvSpPr{cNvSpPr_extra}/>'
        f'<p:nvPr/>'
        f'</p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'<a:prstGeom prst="{prst}"><a:avLst/></a:prstGeom>'
        f'{fill_xml}'
        f'{line_xml}'
        f'</p:spPr>'
        f'{txbody_xml}'
        f'</p:sp>'
    )
    return etree.fromstring(xml_str)


def make_textbox(shape_id, name, x, y, cx, cy, text, font_size, bold, color, align='ctr'):
    """Build a styled text-box <p:sp>."""
    b_str = '1' if bold else '0'
    txbody_xml = (
        f'<p:txBody'
        f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:bodyPr wrap="square"><a:spAutoFit/></a:bodyPr>'
        f'<a:lstStyle/>'
        f'<a:p>'
        f'<a:pPr algn="{align}">'
        f'<a:defRPr sz="{font_size}" b="{b_str}" i="0">'
        f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
        f'</a:defRPr>'
        f'</a:pPr>'
        f'<a:r><a:rPr lang="en-US" dirty="0"/><a:t>{text}</a:t></a:r>'
        f'</a:p>'
        f'</p:txBody>'
    )
    return make_sp(
        shape_id, name, x, y, cx, cy, 'rect',
        no_fill=True, txbody_xml=txbody_xml, is_textbox=True
    )


# ── "Our Mentor" section label ────────────────────────────────────────────────
LABEL_W    = 4000000
LABEL_H    = 400000
LABEL_LEFT = (SLIDE_W - LABEL_W) // 2
LABEL_TOP  = MENTOR_SECTION_Y

sid = next_id(slide)
slide.shapes._spTree.append(
    make_textbox(sid, f'TextBox {sid}',
                 LABEL_LEFT, LABEL_TOP, LABEL_W, LABEL_H,
                 'Our Mentor', 2000, True, COLOR_HEADER, align='ctr')
)

# ── 1. Shadow rounded rect ────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_sp(sid, f'Rounded Rectangle {sid}',
            MENTOR_SHADOW_LEFT, MENTOR_CARD_SHADOW_TOP, CARD_W, CARD_H,
            'roundRect', fill_color=COLOR_CARD_SHADOW, no_line=True)
)

# ── 2. Main rounded rect (white) ──────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_sp(sid, f'Rounded Rectangle {sid}',
            MENTOR_CARD_LEFT, MENTOR_CARD_TOP, CARD_W, CARD_H,
            'roundRect', fill_color=COLOR_CARD_BG, no_line=True)
)

# ── 3. Header rectangle (teal top band) ───────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_sp(sid, f'Rectangle {sid}',
            MENTOR_CARD_LEFT, MENTOR_CARD_TOP, CARD_W, RECT_H,
            'rect', fill_color=COLOR_HEADER, no_line=True)
)

# ── 4. Avatar oval ────────────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_sp(sid, f'Oval {sid}',
            OVAL_LEFT, OVAL_TOP, OVAL_W, OVAL_H,
            'ellipse', fill_color=COLOR_CARD_BG,
            line_w=31750, line_color=COLOR_OVAL_BORDER)
)

# ── 5. Initials textbox ───────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_textbox(sid, f'TextBox {sid}',
                 INITIALS_TB_LEFT, INITIALS_TB_TOP, OVAL_W, OVAL_H,
                 'AKM', 2200, True, COLOR_INITIALS, align='ctr')
)

# ── 6. Full name textbox ──────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_textbox(sid, f'TextBox {sid}',
                 NAME_TB_LEFT, NAME_TB_TOP, 2286000, 457200,
                 'Ashish Kumar Mahato', 1700, True, COLOR_NAME, align='ctr')
)

# ── 7. Divider bar ───────────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_sp(sid, f'Rectangle {sid}',
            DIVIDER_LEFT, DIVIDER_TOP, DIVIDER_W, DIVIDER_H,
            'rect', fill_color=COLOR_DIVIDER, no_line=True)
)

# ── 8. Role label textbox ────────────────────────────────────────────────────
sid = next_id(slide)
slide.shapes._spTree.append(
    make_textbox(sid, f'TextBox {sid}',
                 ROLE_TB_LEFT, ROLE_TB_TOP, 2286000, 731520,
                 'Mentor', 1400, False, COLOR_ROLE, align='ctr')
)

# ── Save ──────────────────────────────────────────────────────────────────────
p.save(OUTPUT_PATH)
print('Done! Mentor card for Ashish Kumar Mahato added to slide 2.')
